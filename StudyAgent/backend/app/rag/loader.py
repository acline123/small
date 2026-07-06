from langchain_community.document_loaders import Docx2txtLoader, PyPDFLoader, TextLoader
from langchain_core.documents import Document


def _load_pptx(file_path: str) -> list[Document]:
    """使用 python-pptx 按幻灯片提取文本。"""
    from pptx import Presentation

    prs = Presentation(file_path)
    documents = []
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            documents.append(
                Document(
                    page_content="\n".join(texts),
                    metadata={"slide_number": slide_num},
                )
            )
    return documents


def load_document(file_path: str, file_type: str):
    """根据文件类型加载文档。"""
    if file_type == "pdf":
        loader = PyPDFLoader(file_path)
        return loader.load()
    elif file_type == "docx":
        loader = Docx2txtLoader(file_path)
        return loader.load()
    elif file_type == "txt":
        loader = TextLoader(file_path, encoding="utf-8")
        return loader.load()
    elif file_type == "pptx":
        return _load_pptx(file_path)
    else:
        raise ValueError(f"不支持的文件类型: {file_type}")
